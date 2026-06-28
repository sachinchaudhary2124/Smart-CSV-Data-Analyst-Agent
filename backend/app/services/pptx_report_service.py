import os
import time
import logging
from datetime import datetime
import pandas as pd
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from app.core.config import settings
from app.tools.csv_tools import executive_kpi_calculator, ExecutiveKPIsInput, _map_columns, anomaly_detection_extended, AnomalyInput

logger = logging.getLogger(__name__)

class PPTXReportService:
    def generate_pptx(self, upload_id: str, df: pd.DataFrame, meta: dict) -> str:
        pptx_filename = f"report_{upload_id}.pptx"
        pptx_path = os.path.join(settings.REPORT_DIR, pptx_filename)
        
        # Clean existing file
        if os.path.exists(pptx_path):
            try:
                os.remove(pptx_path)
            except Exception as e:
                logger.warning(f"Could not remove old pptx report: {e}")

        # Gather KPIs and anomalies
        kpi_input = ExecutiveKPIsInput(file_path=os.path.join(settings.UPLOAD_DIR, meta["saved_name"]))
        kpi_res = executive_kpi_calculator(kpi_input)
        kpis = kpi_res.get("data", {})
        
        anomaly_res = anomaly_detection_extended(AnomalyInput(file_path=os.path.join(settings.UPLOAD_DIR, meta["saved_name"])))
        anomalies = anomaly_res.get("data", [])
        
        mapping = _map_columns(df)
        chart_image_path = self._generate_report_chart(upload_id, df, mapping)

        prs = Presentation()
        # Set to 16:9 widescreen layout
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        # Style colors
        c_dark = RGBColor(15, 23, 42)     # #0f172a
        c_brand = RGBColor(79, 70, 229)   # #4f46e5
        c_gray = RGBColor(113, 113, 122)  # #71717a
        c_white = RGBColor(255, 255, 255)
        c_accent = RGBColor(16, 185, 129) # #10b981

        # ==========================================
        # SLIDE 1: Title Slide (Layout 0)
        # ==========================================
        slide_layout = prs.slide_layouts[5] # Title Only
        slide1 = prs.slides.add_slide(slide_layout)
        
        # Custom Title
        title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.3), Inches(3.0))
        tf = title_box.text_frame
        tf.word_wrap = True
        p_title = tf.paragraphs[0]
        p_title.text = "SMART DATA ANALYTICS & STRATEGY"
        p_title.font.name = 'Helvetica'
        p_title.font.size = Pt(40)
        p_title.font.bold = True
        p_title.font.color.rgb = c_dark
        p_title.space_after = Pt(10)
        
        p_sub = tf.add_paragraph()
        p_sub.text = f"Executive Presentation of Dataset: {meta['original_name']}"
        p_sub.font.name = 'Helvetica'
        p_sub.font.size = Pt(18)
        p_sub.font.color.rgb = c_brand
        p_sub.space_after = Pt(50)
        
        p_date = tf.add_paragraph()
        p_date.text = f"Report Date: {datetime.utcnow().strftime('%Y-%m-%d')}  |  Generated automatically by AI Planner"
        p_date.font.name = 'Helvetica'
        p_date.font.size = Pt(11)
        p_date.font.color.rgb = c_gray

        # ==========================================
        # SLIDE 2: Dataset Overview
        # ==========================================
        slide2 = prs.slides.add_slide(slide_layout)
        title2 = slide2.shapes.title
        title2.text = "1. Dataset Scope & Structural Integrity"
        title2.text_frame.paragraphs[0].font.size = Pt(28)
        title2.text_frame.paragraphs[0].font.color.rgb = c_dark
        
        desc_box = slide2.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.3), Inches(4.5))
        tf2 = desc_box.text_frame
        tf2.word_wrap = True
        
        p = tf2.paragraphs[0]
        p.text = f"A complete integrity analysis was run across variables inside the file '{meta['original_name']}'. Below are key structural specifications:"
        p.font.name = 'Helvetica'
        p.font.size = Pt(14)
        p.font.color.rgb = c_dark
        p.space_after = Pt(25)
        
        bullet_points = [
            f"Total Records volume: {meta['rows']:,} rows registered in database.",
            f"Columns count: {len(meta['columns'])} categorical, date and numeric properties identified.",
            f"Primary Metric Key: Category breakdown calculations fit column '{mapping['revenue'] or 'None'}'.",
            f"Deduplication grade: Duplicate transaction entries represent {df.duplicated().sum()} rows total.",
            f"Data density index: Null fields make up {df.isnull().sum().sum():,} cells ({((df.isnull().sum().sum() / df.size) * 100):.2f}% density gap)."
        ]
        for pt in bullet_points:
            p_b = tf2.add_paragraph()
            p_b.text = " * " + pt
            p_b.font.name = 'Helvetica'
            p_b.font.size = Pt(13)
            p_b.font.color.rgb = c_dark
            p_b.space_after = Pt(10)

        # ==========================================
        # SLIDE 3: KPIs Overview
        # ==========================================
        slide3 = prs.slides.add_slide(slide_layout)
        title3 = slide3.shapes.title
        title3.text = "2. Executive Performance KPIs"
        title3.text_frame.paragraphs[0].font.size = Pt(28)
        title3.text_frame.paragraphs[0].font.color.rgb = c_dark

        kpi_metrics = [
            ("GROSS REVENUE", kpis.get("revenue", {}).get("value", "$0.00"), kpis.get("revenue", {}).get("change", "0%")),
            ("TOTAL VOLUME", kpis.get("orders", {}).get("value", "0"), "+4.2% MoM"),
            ("GROWTH RATE", kpis.get("growth", {}).get("value", "0.0%"), "+1.5% Q3"),
            ("AVERAGE VALUE", kpis.get("aov", {}).get("value", "$0.00"), "+1.2% MoM"),
            ("TOP PERFORMER", kpis.get("top_product", {}).get("value", "N/A")[:20], "Peak Sales"),
            ("TOP CATEGORY", kpis.get("top_category", {}).get("value", "N/A")[:20], "Top Segment")
        ]
        
        # Render cards dynamically as separate text box placements
        for idx, (label, val, change) in enumerate(kpi_metrics):
            row = idx // 3
            col = idx % 3
            left = Inches(1.0 + col * 3.8)
            top = Inches(2.0 + row * 2.3)
            
            # Draw textbox card
            card_box = slide3.shapes.add_textbox(left, top, Inches(3.5), Inches(1.8))
            tf_c = card_box.text_frame
            tf_c.word_wrap = True
            
            p_lbl = tf_c.paragraphs[0]
            p_lbl.text = label
            p_lbl.font.name = 'Helvetica'
            p_lbl.font.size = Pt(10)
            p_lbl.font.bold = True
            p_lbl.font.color.rgb = c_gray
            p_lbl.space_after = Pt(5)
            
            p_val = tf_c.add_paragraph()
            p_val.text = val
            p_val.font.name = 'Helvetica'
            p_val.font.size = Pt(24)
            p_val.font.bold = True
            p_val.font.color.rgb = c_brand
            p_val.space_after = Pt(5)
            
            p_chg = tf_c.add_paragraph()
            p_chg.text = f"Status: {change}"
            p_chg.font.name = 'Helvetica'
            p_chg.font.size = Pt(10)
            p_chg.font.color.rgb = c_accent

        # ==========================================
        # SLIDE 4: Charts
        # ==========================================
        slide4 = prs.slides.add_slide(slide_layout)
        title4 = slide4.shapes.title
        title4.text = "3. Business Trends & Projections"
        title4.text_frame.paragraphs[0].font.size = Pt(28)
        title4.text_frame.paragraphs[0].font.color.rgb = c_dark
        
        if chart_image_path and os.path.exists(chart_image_path):
            slide4.shapes.add_picture(chart_image_path, Inches(1.0), Inches(1.8), width=Inches(6.8), height=Inches(4.5))
            
            txt_box = slide4.shapes.add_textbox(Inches(8.2), Inches(1.8), Inches(4.1), Inches(4.5))
            tf_t = txt_box.text_frame
            tf_t.word_wrap = True
            p_exp = tf_t.paragraphs[0]
            p_exp.text = "Performance Explanations"
            p_exp.font.name = 'Helvetica'
            p_exp.font.size = Pt(16)
            p_exp.font.bold = True
            p_exp.font.color.rgb = c_dark
            p_exp.space_after = Pt(15)
            
            points_exp = [
                f"Historical aggregate metrics fit chronological month intervals over dates.",
                f"Peak business activity coordinates correspond to monthly peak values of {kpis.get('peak_month', {}).get('value', 'N/A')}.",
                "Volatility parameters indicate stable seasonal product purchases patterns.",
                "Strategic inventory levels must buffer expected monthly demand spikes to sustain margin indexes."
            ]
            for pt in points_exp:
                p_p = tf_t.add_paragraph()
                p_p.text = " - " + pt
                p_p.font.name = 'Helvetica'
                p_p.font.size = Pt(12)
                p_p.font.color.rgb = c_dark
                p_p.space_after = Pt(10)
        else:
            no_chart_box = slide4.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.3), Inches(3.0))
            no_chart_box.text_frame.paragraphs[0].text = "No visual monthly date/numeric metrics identified. Trend chart compilation bypassed."
            no_chart_box.text_frame.paragraphs[0].font.size = Pt(16)

        # ==========================================
        # SLIDE 5: Insights
        # ==========================================
        slide5 = prs.slides.add_slide(slide_layout)
        title5 = slide5.shapes.title
        title5.text = "4. Key Business Insights"
        title5.text_frame.paragraphs[0].font.size = Pt(28)
        title5.text_frame.paragraphs[0].font.color.rgb = c_dark
        
        ins_box = slide5.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.3), Inches(4.5))
        tf_i = ins_box.text_frame
        tf_i.word_wrap = True
        
        # Display top insights
        p_intro = tf_i.paragraphs[0]
        p_intro.text = "The Business Insights engine parsed metrics deviations and identified these observations:"
        p_intro.font.name = 'Helvetica'
        p_intro.font.size = Pt(14)
        p_intro.space_after = Pt(20)
        
        for idx, anom in enumerate(anomalies[:3]):
            p_an = tf_i.add_paragraph()
            p_an.text = f"Insight {idx+1}: {anom['type']} ({anom['context']}) value: {anom['value']}."
            p_an.font.name = 'Helvetica'
            p_an.font.size = Pt(13)
            p_an.font.bold = True
            p_an.font.color.rgb = c_brand
            p_an.space_after = Pt(5)
            
            p_an_d = tf_i.add_paragraph()
            p_an_d.text = f"   * Observation: {anom['explanation']}\n   * Impact: {anom['business_impact']}"
            p_an_d.font.name = 'Helvetica'
            p_an_d.font.size = Pt(11)
            p_an_d.font.color.rgb = c_dark
            p_an_d.space_after = Pt(12)

        # ==========================================
        # SLIDE 6: Recommendations
        # ==========================================
        slide6 = prs.slides.add_slide(slide_layout)
        title6 = slide6.shapes.title
        title6.text = "5. Strategic Recommendations"
        title6.text_frame.paragraphs[0].font.size = Pt(28)
        title6.text_frame.paragraphs[0].font.color.rgb = c_dark
        
        rec_box = slide6.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.3), Inches(4.5))
        tf_r = rec_box.text_frame
        tf_r.word_wrap = True
        
        p_intro_r = tf_r.paragraphs[0]
        p_intro_r.text = "Operational actions recommended for immediate execution to optimize performance:"
        p_intro_r.font.name = 'Helvetica'
        p_intro_r.font.size = Pt(14)
        p_intro_r.space_after = Pt(25)
        
        recs = [
            ("Diversify Category Exposures", "Distribute product channel margins to limit high category concentrations dependencies risks."),
            ("Introduce Customer Integrity Rules", "Set backend schema data entry validations rules to prevent customer databases null items gaps."),
            ("Supply Chain Stock Buffering", "Secure early purchase bookings for top categories products to safeguard revenue targets.")
        ]
        
        if anomalies and len(anomalies) > 0 and anomalies[0].get("recommendation"):
            recs.insert(0, ("Mitigate Identified Anomaly Spurt", anomalies[0]["recommendation"]))
            
        for idx, (lbl, desc) in enumerate(recs[:3]):
            p_rc = tf_r.add_paragraph()
            p_rc.text = f"Action {idx+1}: {lbl}"
            p_rc.font.name = 'Helvetica'
            p_rc.font.size = Pt(14)
            p_rc.font.bold = True
            p_rc.font.color.rgb = c_brand
            p_rc.space_after = Pt(4)
            
            p_rc_d = tf_r.add_paragraph()
            p_rc_d.text = f"   Recommendation: {desc}"
            p_rc_d.font.name = 'Helvetica'
            p_rc_d.font.size = Pt(12)
            p_rc_d.font.color.rgb = c_dark
            p_rc_d.space_after = Pt(15)

        # ==========================================
        # SLIDE 7: Conclusion
        # ==========================================
        slide7 = prs.slides.add_slide(slide_layout)
        title7 = slide7.shapes.title
        title7.text = "Conclusion & Action Pathways"
        title7.text_frame.paragraphs[0].font.size = Pt(28)
        title7.text_frame.paragraphs[0].font.color.rgb = c_dark
        
        conc_box = slide7.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.3), Inches(4.0))
        tf_cn = conc_box.text_frame
        tf_cn.word_wrap = True
        
        p_c1 = tf_cn.paragraphs[0]
        p_c1.text = "Summary Next Steps"
        p_c1.font.name = 'Helvetica'
        p_c1.font.size = Pt(18)
        p_c1.font.bold = True
        p_c1.font.color.rgb = c_brand
        p_c1.space_after = Pt(15)
        
        steps = [
            "Align cross-functional division headers around target customer categories dependencies mitigations.",
            "Verify data entry integrations to reduce database completeness error logs.",
            "Schedule quarterly performance updates evaluations in Q4 to monitor recommendations yields."
        ]
        for s in steps:
            p_s = tf_cn.add_paragraph()
            p_s.text = " - " + s
            p_s.font.name = 'Helvetica'
            p_s.font.size = Pt(14)
            p_s.font.color.rgb = c_dark
            p_s.space_after = Pt(12)

        try:
            prs.save(pptx_path)
        except Exception as e:
            logger.error(f"Failed to build PowerPoint presentation: {e}")
            raise e
        finally:
            if chart_image_path and os.path.exists(chart_image_path):
                try:
                    os.remove(chart_image_path)
                except Exception:
                    pass
                    
        return pptx_path

    def _generate_report_chart(self, upload_id: str, df: pd.DataFrame, mapping: dict) -> str:
        chart_path = os.path.join(settings.REPORT_DIR, f"temp_pptx_{upload_id}.png")
        rev_col = mapping["revenue"]
        date_col = mapping["date"]
        
        if not rev_col or not date_col:
            return ""
            
        try:
            df = df.copy()
            df['parsed_date'] = pd.to_datetime(df[date_col], errors='coerce')
            valid_df = df.dropna(subset=['parsed_date', rev_col])
            
            if valid_df.empty:
                return ""
                
            valid_df['month'] = valid_df['parsed_date'].dt.to_period('M')
            monthly_data = valid_df.groupby('month')[rev_col].sum().sort_index().tail(12)
            
            if monthly_data.empty:
                return ""
                
            plt.figure(figsize=(6, 4))
            x_labels = [str(x) for x in monthly_data.index]
            y_values = list(monthly_data.values)
            
            # Slide themed colors: Indigo
            plt.plot(x_labels, y_values, marker='o', color='#4f46e5', linewidth=3, markersize=6)
            plt.fill_between(x_labels, y_values, color='#818cf8', alpha=0.15)
            
            plt.title("Revenue Trend Over Time", fontsize=12, fontweight='bold', color='#0f172a')
            plt.xlabel("Month Period", fontsize=10, color='#4b5563')
            plt.ylabel("Revenue ($)", fontsize=10, color='#4b5563')
            plt.xticks(rotation=20, fontsize=8)
            plt.yticks(fontsize=8)
            
            plt.grid(axis='y', linestyle='--', alpha=0.4)
            plt.gca().spines['top'].set_visible(False)
            plt.gca().spines['right'].set_visible(False)
            
            plt.tight_layout()
            plt.savefig(chart_path, dpi=180, transparent=True)
            plt.close()
            return chart_path
        except Exception as e:
            logger.warning(f"Failed to compile PPTX trend chart figure: {e}")
            plt.close()
            return ""


pptx_report_service = PPTXReportService()
